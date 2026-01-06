"""
PPTX Renderer V2 - Professional Layout Engine

This module provides visually appealing slide rendering with:
1. Fixed text overlap issues
2. Modern visual elements (accent lines, shapes, gradients)
3. Intent-aware layouts
4. Better spacing and visual hierarchy
5. Decorative elements for polish

Author: SlideGen Team
"""

from typing import Any, List, Dict, Optional
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from .overflow import BoundingBox, TextOverflowEngine
from .themes import ThemeColorScheme


class ShapeFactoryV2:
    """Enhanced shape factory with more visual elements."""
    
    @staticmethod
    def create_text_box(slide, box: BoundingBox, text: str, font_size: float,
                        color, bold: bool = False, align=PP_ALIGN.LEFT,
                        vertical_anchor=MSO_ANCHOR.TOP):
        """Create a text box with proper anchoring."""
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        # Set vertical anchor to prevent overlap
        try:
            tf.anchor = vertical_anchor
        except:
            pass
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        
        # Add line spacing
        p.line_spacing = 1.2
        
        return shape
    
    @staticmethod
    def create_bullet_list(slide, box: BoundingBox, items: List[Dict],
                           font_size: float, color, bullet_color=None,
                           line_spacing: float = 1.5):
        """Create bullet list with proper spacing to prevent overlap."""
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = shape.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            text = item.get('text', '') if isinstance(item, dict) else str(item)
            level = item.get('level', 0) if isinstance(item, dict) else 0
            
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.level = min(level, 1)  # Max 2 levels
            p.bullet = True
            p.line_spacing = line_spacing
            
            # Add space after each paragraph
            p.space_after = Pt(8)
        
        return shape
    
    @staticmethod
    def create_rectangle(slide, box: BoundingBox, color, alpha: float = 1.0):
        """Create a solid rectangle."""
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    @staticmethod
    def create_accent_line(slide, x: float, y: float, width: float, color, thickness: float = 4):
        """Create a horizontal accent line for visual separation."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Pt(thickness)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    @staticmethod
    def create_vertical_line(slide, x: float, y: float, height: float, color, thickness: float = 2):
        """Create a vertical divider line."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Pt(thickness), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    @staticmethod
    def create_rounded_rectangle(slide, box: BoundingBox, color, radius: float = 0.2):
        """Create a rounded rectangle for modern look."""
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    @staticmethod
    def create_circle(slide, x: float, y: float, diameter: float, color):
        """Create a circle shape."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(diameter), Inches(diameter)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    @staticmethod
    def create_number_badge(slide, x: float, y: float, number: int, 
                           bg_color, text_color, size: float = 0.5):
        """Create a numbered badge/circle."""
        # Background circle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        
        # Add number text
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(int(size * 28))
        p.font.color.rgb = text_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        return shape


class SlideRendererV2:
    """
    Enhanced slide renderer with professional layouts.
    
    Features:
    - Intent-aware layouts
    - Visual hierarchy
    - Accent elements
    - Proper spacing to prevent overlap
    """
    
    # Slide dimensions (16:9)
    WIDTH = 13.333
    HEIGHT = 7.5
    MARGIN = 0.6
    
    def __init__(self, theme: ThemeColorScheme):
        self.theme = theme
        self.overflow = TextOverflowEngine()
        self.factory = ShapeFactoryV2()
    
    def get_layout(self, slide_type: str) -> Dict[str, BoundingBox]:
        """Get layout bounding boxes for metrics evaluation."""
        w = self.WIDTH - 2 * self.MARGIN
        layouts = {
            'title': {
                'title': BoundingBox(self.MARGIN, 0.8, w, 1.6),
                'content': BoundingBox(self.MARGIN, 3.4, w, 1.0)
            },
            'content': {
                'title': BoundingBox(self.MARGIN, 0.4, w, 0.9),
                'content': BoundingBox(self.MARGIN, 1.7, w, 5.3)
            },
            'two_column': {
                'title': BoundingBox(self.MARGIN, 0.3, w, 0.8),
                'left': BoundingBox(self.MARGIN, 2.4, w/2 - 0.4, 4.5),
                'right': BoundingBox(self.MARGIN + w/2 + 0.4, 2.4, w/2 - 0.4, 4.5)
            },
            'closing': {
                'title': BoundingBox(self.MARGIN, 2.5, w, 1.5),
                'content': BoundingBox(self.MARGIN, 4.2, w, 0.8)
            },
            'section': {
                'title': BoundingBox(1, 2.5, w - 2, 1.5),
                'content': BoundingBox(1, 4.2, w - 2, 0.8)
            },
        }
        return layouts.get(slide_type, layouts['content'])
    
    def render(self, prs: Presentation, data: Dict) -> Any:
        """Render a slide based on its type and intent."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        slide_type = data.get('slide_type', 'content')
        intent = data.get('intent', '')
        
        # Route to appropriate renderer
        if slide_type == 'title' or intent == 'vision':
            return self._render_title_slide(slide, data)
        elif slide_type == 'section':
            return self._render_section_slide(slide, data)
        elif slide_type == 'two_column' or intent == 'comparison':
            return self._render_two_column_slide(slide, data)
        elif slide_type == 'closing' or intent == 'call_to_action':
            return self._render_closing_slide(slide, data)
        elif intent == 'agenda':
            return self._render_agenda_slide(slide, data)
        elif intent == 'data_insight':
            return self._render_data_slide(slide, data)
        else:
            return self._render_content_slide(slide, data)
    
    def _render_title_slide(self, slide, data: Dict):
        """Render an impactful title/vision slide."""
        w = self.WIDTH
        h = self.HEIGHT
        
        # Full-width colored header bar
        self.factory.create_rectangle(
            slide,
            BoundingBox(0, 0, w, 2.8),
            self.theme.get_rgb('primary')
        )
        
        # Accent line at bottom of header
        self.factory.create_accent_line(
            slide, 0, 2.8, w,
            self.theme.get_rgb('accent'),
            thickness=6
        )
        
        # Small decorative circle
        self.factory.create_circle(
            slide, w - 1.5, 0.4, 0.3,
            self.theme.get_rgb('accent')
        )
        
        # Title
        title = data.get('title', '')
        if title:
            result = self.overflow.fit_text(title, BoundingBox(self.MARGIN, 0.8, w - 2*self.MARGIN, 1.6), 44)
            self.factory.create_text_box(
                slide,
                BoundingBox(self.MARGIN, 0.8, w - 2*self.MARGIN, 1.6),
                result['text'],
                result['font_size'],
                self.theme.get_rgb('text_light'),
                bold=True,
                align=PP_ALIGN.LEFT,
                vertical_anchor=MSO_ANCHOR.MIDDLE
            )
        
        # Subtitle
        subtitle = data.get('subtitle', '')
        if subtitle:
            result = self.overflow.fit_text(subtitle, BoundingBox(self.MARGIN, 3.4, w - 2*self.MARGIN, 1.0), 24)
            self.factory.create_text_box(
                slide,
                BoundingBox(self.MARGIN, 3.4, w - 2*self.MARGIN, 1.0),
                result['text'],
                result['font_size'],
                self.theme.get_rgb('text_dark'),
                bold=False,
                align=PP_ALIGN.LEFT
            )
        
        # Bottom decorative element
        self.factory.create_rectangle(
            slide,
            BoundingBox(0, h - 0.3, 3, 0.3),
            self.theme.get_rgb('secondary')
        )
        
        return slide
    
    def _render_section_slide(self, slide, data: Dict):
        """Render a section divider slide."""
        w = self.WIDTH
        h = self.HEIGHT
        
        # Side accent bar
        self.factory.create_rectangle(
            slide,
            BoundingBox(0, 0, 0.15, h),
            self.theme.get_rgb('primary')
        )
        
        # Large circle decoration
        self.factory.create_circle(
            slide, w - 2.5, h - 2.5, 2,
            self._lighten_color(self.theme.get_rgb('secondary'), 0.8)
        )
        
        # Section title
        title = data.get('title', '')
        if title:
            result = self.overflow.fit_text(title, BoundingBox(1, 2.5, w - 3, 1.5), 40)
            self.factory.create_text_box(
                slide,
                BoundingBox(1, 2.5, w - 3, 1.5),
                result['text'],
                result['font_size'],
                self.theme.get_rgb('primary'),
                bold=True,
                align=PP_ALIGN.LEFT
            )
        
        # Subtitle
        subtitle = data.get('subtitle', '')
        if subtitle:
            self.factory.create_text_box(
                slide,
                BoundingBox(1, 4.2, w - 3, 0.8),
                subtitle[:100],
                20,
                self.theme.get_rgb('text_dark'),
                bold=False,
                align=PP_ALIGN.LEFT
            )
        
        return slide
    
    def _render_content_slide(self, slide, data: Dict):
        """Render a standard content slide with proper spacing."""
        w = self.WIDTH - 2 * self.MARGIN
        
        # Header accent bar
        self.factory.create_rectangle(
            slide,
            BoundingBox(0, 0, 0.1, self.HEIGHT),
            self.theme.get_rgb('primary')
        )
        
        # Title with accent underline
        title = data.get('title', '')
        if title:
            result = self.overflow.fit_text(title, BoundingBox(self.MARGIN, 0.4, w, 0.9), 28)
            self.factory.create_text_box(
                slide,
                BoundingBox(self.MARGIN, 0.4, w, 0.9),
                result['text'],
                result['font_size'],
                self.theme.get_rgb('primary'),
                bold=True,
                align=PP_ALIGN.LEFT
            )
            
            # Accent line under title
            self.factory.create_accent_line(
                slide, self.MARGIN, 1.35, 2,
                self.theme.get_rgb('accent'),
                thickness=4
            )
        
        # Body points - with proper spacing
        points = data.get('body_points', [])
        if points:
            # Filter and limit points
            processed = []
            for p in points[:5]:  # Max 5 points
                text = p.get('text', '') if isinstance(p, dict) else str(p)
                level = p.get('level', 0) if isinstance(p, dict) else 0
                processed.append({'text': text[:80], 'level': level})  # Limit text length
            
            # Content area starts lower to avoid title overlap
            self.factory.create_bullet_list(
                slide,
                BoundingBox(self.MARGIN, 1.7, w, 5.3),
                processed,
                20,
                self.theme.get_rgb('text_dark'),
                line_spacing=1.6
            )
        
        # Small decorative element in corner
        self.factory.create_circle(
            slide, self.WIDTH - 0.8, self.HEIGHT - 0.8, 0.25,
            self._lighten_color(self.theme.get_rgb('accent'), 0.5)
        )
        
        return slide
    
    def _render_two_column_slide(self, slide, data: Dict):
        """Render a comparison/two-column slide."""
        w = self.WIDTH - 2 * self.MARGIN
        col_width = (w - 0.8) / 2  # Gap between columns
        
        # Header bar
        self.factory.create_rectangle(
            slide,
            BoundingBox(0, 0, self.WIDTH, 1.4),
            self.theme.get_rgb('primary')
        )
        
        # Title
        title = data.get('title', '')
        if title:
            result = self.overflow.fit_text(title, BoundingBox(self.MARGIN, 0.3, w, 0.8), 26)
            self.factory.create_text_box(
                slide,
                BoundingBox(self.MARGIN, 0.3, w, 0.8),
                result['text'],
                result['font_size'],
                self.theme.get_rgb('text_light'),
                bold=True,
                align=PP_ALIGN.LEFT
            )
        
        # Column headers
        left_header = data.get('left_header', 'Option A')
        right_header = data.get('right_header', 'Option B')
        
        # Left column header with background
        self.factory.create_rounded_rectangle(
            slide,
            BoundingBox(self.MARGIN, 1.7, col_width, 0.5),
            self._lighten_color(self.theme.get_rgb('secondary'), 0.3)
        )
        self.factory.create_text_box(
            slide,
            BoundingBox(self.MARGIN + 0.2, 1.75, col_width - 0.4, 0.4),
            left_header[:30],
            16,
            self.theme.get_rgb('primary'),
            bold=True,
            align=PP_ALIGN.CENTER
        )
        
        # Right column header with background
        self.factory.create_rounded_rectangle(
            slide,
            BoundingBox(self.MARGIN + col_width + 0.8, 1.7, col_width, 0.5),
            self._lighten_color(self.theme.get_rgb('accent'), 0.3)
        )
        self.factory.create_text_box(
            slide,
            BoundingBox(self.MARGIN + col_width + 1, 1.75, col_width - 0.4, 0.4),
            right_header[:30],
            16,
            self.theme.get_rgb('primary'),
            bold=True,
            align=PP_ALIGN.CENTER
        )
        
        # Vertical divider
        self.factory.create_vertical_line(
            slide,
            self.MARGIN + col_width + 0.35, 2.4, 4.5,
            self._lighten_color(self.theme.get_rgb('text_dark'), 0.7),
            thickness=2
        )
        
        # Left column content
        left_items = data.get('left_column', [])
        if left_items:
            processed = self._process_column_items(left_items)
            self.factory.create_bullet_list(
                slide,
                BoundingBox(self.MARGIN, 2.4, col_width, 4.5),
                processed,
                18,
                self.theme.get_rgb('text_dark'),
                line_spacing=1.8
            )
        
        # Right column content
        right_items = data.get('right_column', [])
        if right_items:
            processed = self._process_column_items(right_items)
            self.factory.create_bullet_list(
                slide,
                BoundingBox(self.MARGIN + col_width + 0.8, 2.4, col_width, 4.5),
                processed,
                18,
                self.theme.get_rgb('text_dark'),
                line_spacing=1.8
            )
        
        return slide
    
    def _render_agenda_slide(self, slide, data: Dict):
        """Render a numbered agenda slide."""
        w = self.WIDTH - 2 * self.MARGIN
        
        # Header
        self.factory.create_rectangle(
            slide,
            BoundingBox(0, 0, self.WIDTH, 1.4),
            self.theme.get_rgb('primary')
        )
        
        # Title
        title = data.get('title', 'Agenda')
        self.factory.create_text_box(
            slide,
            BoundingBox(self.MARGIN, 0.35, w, 0.7),
            title[:50],
            28,
            self.theme.get_rgb('text_light'),
            bold=True,
            align=PP_ALIGN.LEFT
        )
        
        # Numbered items with badges
        points = data.get('body_points', [])
        y_start = 1.9
        y_spacing = 1.0
        
        for i, point in enumerate(points[:5]):
            text = point.get('text', '') if isinstance(point, dict) else str(point)
            
            # Number badge
            self.factory.create_number_badge(
                slide,
                self.MARGIN, y_start + i * y_spacing,
                i + 1,
                self.theme.get_rgb('accent'),
                self.theme.get_rgb('text_light'),
                size=0.45
            )
            
            # Item text
            self.factory.create_text_box(
                slide,
                BoundingBox(self.MARGIN + 0.7, y_start + i * y_spacing + 0.05, w - 0.7, 0.5),
                text[:60],
                20,
                self.theme.get_rgb('text_dark'),
                bold=False,
                align=PP_ALIGN.LEFT
            )
            
            # Subtle separator line
            if i < len(points) - 1:
                self.factory.create_accent_line(
                    slide,
                    self.MARGIN + 0.7,
                    y_start + (i + 1) * y_spacing - 0.15,
                    w - 0.7,
                    self._lighten_color(self.theme.get_rgb('text_dark'), 0.85),
                    thickness=1
                )
        
        return slide
    
    def _render_data_slide(self, slide, data: Dict):
        """Render a data/insight focused slide with emphasis on numbers."""
        w = self.WIDTH - 2 * self.MARGIN
        
        # Accent bar
        self.factory.create_rectangle(
            slide,
            BoundingBox(0, 0, 0.15, self.HEIGHT),
            self.theme.get_rgb('accent')
        )
        
        # Title
        title = data.get('title', '')
        if title:
            self.factory.create_text_box(
                slide,
                BoundingBox(self.MARGIN, 0.4, w, 1.0),
                title[:80],
                26,
                self.theme.get_rgb('primary'),
                bold=True,
                align=PP_ALIGN.LEFT
            )
        
        # Data points with emphasis styling
        points = data.get('body_points', [])
        y_start = 1.8
        y_spacing = 1.4
        
        for i, point in enumerate(points[:4]):
            text = point.get('text', '') if isinstance(point, dict) else str(point)
            
            # Highlight box for each data point
            self.factory.create_rounded_rectangle(
                slide,
                BoundingBox(self.MARGIN, y_start + i * y_spacing, w, 1.1),
                self._lighten_color(self.theme.get_rgb('secondary'), 0.85)
            )
            
            # Point text
            self.factory.create_text_box(
                slide,
                BoundingBox(self.MARGIN + 0.3, y_start + i * y_spacing + 0.25, w - 0.6, 0.6),
                text[:80],
                20,
                self.theme.get_rgb('text_dark'),
                bold=False,
                align=PP_ALIGN.LEFT
            )
        
        return slide
    
    def _render_closing_slide(self, slide, data: Dict):
        """Render a closing/thank you slide."""
        w = self.WIDTH
        h = self.HEIGHT
        
        # Full background color
        self.factory.create_rectangle(
            slide,
            BoundingBox(0, 0, w, h),
            self.theme.get_rgb('primary')
        )
        
        # Decorative circles
        self.factory.create_circle(
            slide, -0.5, h - 2, 2.5,
            self._lighten_color(self.theme.get_rgb('secondary'), 0.3)
        )
        self.factory.create_circle(
            slide, w - 1.5, -0.5, 2,
            self._lighten_color(self.theme.get_rgb('accent'), 0.4)
        )
        
        # Title
        title = data.get('title', 'Thank You')
        self.factory.create_text_box(
            slide,
            BoundingBox(self.MARGIN, 2.5, w - 2*self.MARGIN, 1.5),
            title[:50],
            48,
            self.theme.get_rgb('text_light'),
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE
        )
        
        # Subtitle
        subtitle = data.get('subtitle', '')
        if subtitle:
            self.factory.create_text_box(
                slide,
                BoundingBox(self.MARGIN, 4.2, w - 2*self.MARGIN, 0.8),
                subtitle[:80],
                24,
                self._lighten_color(self.theme.get_rgb('text_light'), 0.2),
                bold=False,
                align=PP_ALIGN.CENTER
            )
        
        # Bottom accent line
        self.factory.create_accent_line(
            slide,
            (w - 3) / 2, 5.3, 3,
            self.theme.get_rgb('accent'),
            thickness=4
        )
        
        return slide
    
    def _process_column_items(self, items: List) -> List[Dict]:
        """Process column items for consistent format."""
        processed = []
        for item in items[:4]:  # Max 4 items per column
            if isinstance(item, dict):
                text = item.get('text', '')[:60]
                level = item.get('level', 0)
            else:
                text = str(item)[:60]
                level = 0
            processed.append({'text': text, 'level': level})
        return processed
    
    def _lighten_color(self, rgb_color: RGBColor, factor: float) -> RGBColor:
        """Lighten a color by mixing with white."""
        r = int(rgb_color[0] + (255 - rgb_color[0]) * factor)
        g = int(rgb_color[1] + (255 - rgb_color[1]) * factor)
        b = int(rgb_color[2] + (255 - rgb_color[2]) * factor)
        return RGBColor(min(r, 255), min(g, 255), min(b, 255))


__all__ = ['ShapeFactoryV2', 'SlideRendererV2']

