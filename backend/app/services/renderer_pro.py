"""
PPTX Renderer Pro - Professional Meeting-Ready Layouts

This module provides polished, business-ready slide rendering with:
1. Rich typography system (varied fonts, weights, colors)
2. Vector graphics icons (shapes-based, not emoji)
3. Multiple layout variations (cards, grids, timelines)
4. Professional color usage and gradients
5. Optional web image integration

Suitable for: Team meetings, project reviews, corporate presentations

Author: SlideGen Team
"""

from typing import Any, List, Dict, Optional, Tuple
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap
from .overflow import BoundingBox, TextOverflowEngine
from .themes import ThemeColorScheme
from .smart_typography import smart_typography, adaptive_sizer, TYPOGRAPHY_PRESETS
import os
import requests
from io import BytesIO


# =============================================================================
# TYPOGRAPHY SYSTEM
# =============================================================================

class Typography:
    """Professional typography presets."""
    
    # Font stacks (fallbacks for cross-platform)
    HEADING = "Calibri Light"
    BODY = "Calibri"
    ACCENT = "Calibri"
    MONO = "Consolas"
    
    # Size scale
    SIZES = {
        'hero': 48,
        'h1': 36,
        'h2': 28,
        'h3': 22,
        'h4': 18,
        'body': 16,
        'body_large': 18,
        'caption': 12,
        'small': 10,
        'tiny': 9,
    }
    
    # Weight presets (bold or not, since python-pptx limited)
    WEIGHTS = {
        'light': False,
        'regular': False,
        'medium': False,  # Will use spacing instead
        'bold': True,
        'heavy': True,
    }


# =============================================================================
# VECTOR ICON SYSTEM (Shape-based, not emoji)
# =============================================================================

class VectorIcons:
    """Create professional vector icons using shape combinations."""
    
    @staticmethod
    def create_lightbulb(slide, x: float, y: float, size: float, color: RGBColor):
        """Create a lightbulb icon using shapes."""
        # Bulb (oval)
        bulb = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + size*0.15), Inches(y),
            Inches(size*0.7), Inches(size*0.6)
        )
        bulb.fill.solid()
        bulb.fill.fore_color.rgb = color
        bulb.line.fill.background()
        
        # Base (rectangle)
        base = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + size*0.3), Inches(y + size*0.55),
            Inches(size*0.4), Inches(size*0.25)
        )
        base.fill.solid()
        base.fill.fore_color.rgb = color
        base.line.fill.background()
        
        return bulb
    
    @staticmethod
    def create_chart_bars(slide, x: float, y: float, size: float, colors: List[RGBColor]):
        """Create a bar chart icon."""
        bar_width = size * 0.22
        gap = size * 0.08
        heights = [0.5, 0.8, 0.6, 1.0]  # Relative heights
        
        for i, h in enumerate(heights):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + i * (bar_width + gap)),
                Inches(y + size * (1 - h)),
                Inches(bar_width),
                Inches(size * h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = colors[i % len(colors)]
            bar.line.fill.background()
        
        return bar
    
    @staticmethod
    def create_target(slide, x: float, y: float, size: float, colors: List[RGBColor]):
        """Create a target/bullseye icon."""
        # Outer ring
        for i, ratio in enumerate([1.0, 0.7, 0.4]):
            ring = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + size*(1-ratio)/2), Inches(y + size*(1-ratio)/2),
                Inches(size*ratio), Inches(size*ratio)
            )
            ring.fill.solid()
            ring.fill.fore_color.rgb = colors[i % len(colors)]
            ring.line.fill.background()
        
        return ring
    
    @staticmethod
    def create_arrow_up(slide, x: float, y: float, size: float, color: RGBColor):
        """Create an upward arrow icon."""
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.UP_ARROW,
            Inches(x), Inches(y),
            Inches(size*0.6), Inches(size)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()
        return arrow
    
    @staticmethod
    def create_gear(slide, x: float, y: float, size: float, color: RGBColor):
        """Create a gear/cog icon using octagon."""
        gear = slide.shapes.add_shape(
            MSO_SHAPE.OCTAGON,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        gear.fill.solid()
        gear.fill.fore_color.rgb = color
        gear.line.fill.background()
        
        # Center hole
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + size*0.3), Inches(y + size*0.3),
            Inches(size*0.4), Inches(size*0.4)
        )
        center.fill.solid()
        center.fill.fore_color.rgb = RGBColor(255, 255, 255)
        center.line.fill.background()
        
        return gear
    
    @staticmethod
    def create_checkmark(slide, x: float, y: float, size: float, color: RGBColor):
        """Create a checkmark in circle."""
        # Circle background
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        # Checkmark text (using special character)
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(int(size * 36))
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        return circle
    
    @staticmethod  
    def create_document(slide, x: float, y: float, size: float, color: RGBColor):
        """Create a document/page icon."""
        # Main page
        doc = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + size*0.1), Inches(y),
            Inches(size*0.7), Inches(size)
        )
        doc.fill.solid()
        doc.fill.fore_color.rgb = color
        doc.line.fill.background()
        
        # Lines on document
        for i in range(3):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + size*0.2), Inches(y + size*(0.25 + i*0.2)),
                Inches(size*0.5), Pt(2)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(255, 255, 255)
            line.line.fill.background()
        
        return doc
    
    @staticmethod
    def create_users(slide, x: float, y: float, size: float, color: RGBColor):
        """Create a people/users icon."""
        # Person 1 (head)
        head1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + size*0.25), Inches(y),
            Inches(size*0.3), Inches(size*0.3)
        )
        head1.fill.solid()
        head1.fill.fore_color.rgb = color
        head1.line.fill.background()
        
        # Person 1 (body)
        body1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + size*0.1), Inches(y + size*0.35),
            Inches(size*0.6), Inches(size*0.5)
        )
        body1.fill.solid()
        body1.fill.fore_color.rgb = color
        body1.line.fill.background()
        
        # Person 2 (smaller, behind)
        head2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + size*0.5), Inches(y + size*0.1),
            Inches(size*0.25), Inches(size*0.25)
        )
        head2.fill.solid()
        head2.fill.fore_color.rgb = VectorIcons._lighten(color, 0.3)
        head2.line.fill.background()
        
        return head1
    
    @staticmethod
    def _lighten(color: RGBColor, factor: float) -> RGBColor:
        r = int(color[0] + (255 - color[0]) * factor)
        g = int(color[1] + (255 - color[1]) * factor)
        b = int(color[2] + (255 - color[2]) * factor)
        return RGBColor(min(r, 255), min(g, 255), min(b, 255))


# Intent to vector icon mapping
INTENT_ICON_CREATORS = {
    'vision': 'target',
    'concept_overview': 'lightbulb',
    'framework': 'gear',
    'comparison': 'chart_bars',
    'case_example': 'document',
    'data_insight': 'chart_bars',
    'implications': 'arrow_up',
    'risks_challenges': 'target',
    'future_directions': 'arrow_up',
    'summary_takeaways': 'checkmark',
    'call_to_action': 'target',
    'agenda': 'document',
}


# =============================================================================
# WEB IMAGE INTEGRATION (Optional)
# =============================================================================

class WebImageFetcher:
    """Fetch relevant images from free image APIs."""
    
    # Unsplash API (free tier)
    UNSPLASH_ACCESS_KEY = os.getenv('UNSPLASH_ACCESS_KEY', '')
    
    # Fallback: use placeholder service
    PLACEHOLDER_URL = "https://picsum.photos"
    
    @classmethod
    def fetch_image(cls, query: str, width: int = 400, height: int = 300) -> Optional[BytesIO]:
        """
        Fetch an image related to the query.
        Returns BytesIO if successful, None otherwise.
        """
        try:
            # Try Unsplash first if API key available
            if cls.UNSPLASH_ACCESS_KEY:
                return cls._fetch_unsplash(query, width, height)
            else:
                # Use placeholder service
                return cls._fetch_placeholder(width, height)
        except Exception as e:
            print(f"Image fetch failed: {e}")
            return None
    
    @classmethod
    def _fetch_unsplash(cls, query: str, width: int, height: int) -> Optional[BytesIO]:
        """Fetch from Unsplash API."""
        url = f"https://api.unsplash.com/photos/random"
        params = {
            'query': query,
            'w': width,
            'h': height,
            'client_id': cls.UNSPLASH_ACCESS_KEY
        }
        response = requests.get(url, params=params, timeout=5)
        if response.status_code == 200:
            data = response.json()
            img_url = data['urls']['small']
            img_response = requests.get(img_url, timeout=5)
            if img_response.status_code == 200:
                return BytesIO(img_response.content)
        return None
    
    @classmethod
    def _fetch_placeholder(cls, width: int, height: int) -> Optional[BytesIO]:
        """Fetch from placeholder service."""
        url = f"{cls.PLACEHOLDER_URL}/{width}/{height}"
        response = requests.get(url, timeout=5)
        if response.status_code == 200:
            return BytesIO(response.content)
        return None


class ProShapeFactory:
    """Professional shape factory with enhanced visual elements."""
    
    @staticmethod
    def create_text_box(slide, box: BoundingBox, text: str, font_size: float,
                        color, bold: bool = False, align=PP_ALIGN.LEFT,
                        vertical_anchor=MSO_ANCHOR.TOP, italic: bool = False,
                        font_name: str = None, line_spacing: float = 1.15,
                        letter_spacing: float = 0):
        """Create a text box with advanced typography options."""
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = shape.text_frame
        tf.word_wrap = True
        
        try:
            tf.anchor = vertical_anchor
        except:
            pass
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = align
        p.line_spacing = line_spacing
        
        # Apply font
        p.font.name = font_name or Typography.BODY
        
        return shape
    
    @staticmethod
    def create_rich_text(slide, box: BoundingBox, segments: List[Dict],
                         align=PP_ALIGN.LEFT):
        """
        Create text with multiple styles in one box.
        
        segments: [{'text': 'Hello', 'bold': True, 'color': RGBColor, 'size': 18}, ...]
        """
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        
        for i, seg in enumerate(segments):
            if i == 0:
                run = p.runs[0] if p.runs else p.add_run()
            else:
                run = p.add_run()
            
            run.text = seg.get('text', '')
            run.font.size = Pt(seg.get('size', 16))
            run.font.bold = seg.get('bold', False)
            run.font.italic = seg.get('italic', False)
            if 'color' in seg:
                run.font.color.rgb = seg['color']
            run.font.name = seg.get('font', Typography.BODY)
        
        return shape
    
    @staticmethod
    def create_styled_bullet_list(slide, box: BoundingBox, items: List[Dict],
                                  base_font_size: float, text_color, accent_color,
                                  style: str = 'default'):
        """Create bullet list with rich typography styles."""
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = shape.text_frame
        tf.word_wrap = True
        
        # Bullet characters by style
        bullets = {
            'default': ['▸', '▹', '·'],
            'numbered': ['1.', '2.', '3.', '4.', '5.'],
            'check': ['•', '•', '·'],
            'arrow': ['–', '–', '·'],
        }
        bullet_chars = bullets.get(style, bullets['default'])
        
        for i, item in enumerate(items):
            text = item.get('text', '') if isinstance(item, dict) else str(item)
            level = item.get('level', 0) if isinstance(item, dict) else 0
            priority = item.get('priority', 'normal') if isinstance(item, dict) else 'normal'
            
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            
            # Get bullet character
            if style == 'numbered':
                bullet = bullet_chars[i] if i < len(bullet_chars) else f"({i+1})"
            else:
                bullet = bullet_chars[min(level, len(bullet_chars)-1)]
            
            # Apply different styles based on priority
            if priority == 'critical':
                p.text = f"{bullet}  {text}"
                p.font.bold = True
                p.font.color.rgb = accent_color
                p.font.size = Pt(base_font_size + 1)
                p.font.name = Typography.ACCENT
            elif priority == 'high':
                p.text = f"{bullet}  {text}"
                p.font.bold = True
                p.font.color.rgb = text_color
                p.font.size = Pt(base_font_size)
                p.font.name = Typography.BODY
            else:
                indent = "    " if level > 0 else ""
                p.text = f"{indent}{bullet}  {text}"
                p.font.color.rgb = text_color
                p.font.size = Pt(base_font_size - 2 if level > 0 else base_font_size - 1)
                p.font.name = Typography.BODY
            
            p.level = 0
            p.line_spacing = 1.6
            p.space_after = Pt(12)
        
        return shape
    
    @staticmethod
    def create_card(slide, box: BoundingBox, bg_color, border_color=None, 
                    shadow: bool = True):
        """Create a card-style container."""
        left, top, width, height = box.to_emu()
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(left), Emu(top), Emu(width), Emu(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        
        return shape
    
    @staticmethod
    def create_icon_circle(slide, x: float, y: float, icon: str, 
                          bg_color, icon_color, size: float = 0.6):
        """Create a circular icon badge."""
        # Background circle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        
        # Icon text
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(int(size * 24))
        p.font.color.rgb = icon_color
        p.alignment = PP_ALIGN.CENTER
        
        try:
            tf.anchor = MSO_ANCHOR.MIDDLE
        except:
            pass
        
        return shape
    
    @staticmethod
    def create_stat_box(slide, x: float, y: float, width: float, height: float,
                       number: str, label: str, bg_color, number_color, label_color):
        """Create a statistics/metric display box with auto-sizing text."""
        # Background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        
        # Number (large) - auto-size based on length
        num_font_size = 32 if len(number) <= 6 else (26 if len(number) <= 10 else 20)
        num_shape = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.15),
            Inches(width - 0.2), Inches(height * 0.5)
        )
        tf = num_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(num_font_size)
        p.font.bold = True
        p.font.color.rgb = number_color
        p.alignment = PP_ALIGN.CENTER
        
        # Label (small) - auto-size based on length, with word wrap
        label_font_size = 12 if len(label) <= 30 else (10 if len(label) <= 50 else 9)
        label_shape = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + height * 0.5),
            Inches(width - 0.2), Inches(height * 0.45)
        )
        tf = label_shape.text_frame
        tf.word_wrap = True  # Enable word wrap for labels
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(label_font_size)
        p.font.color.rgb = label_color
        p.alignment = PP_ALIGN.CENTER
        
        return shape
    
    @staticmethod
    def create_timeline_node(slide, x: float, y: float, 
                            text: str, is_active: bool,
                            primary_color, secondary_color, text_color):
        """Create a timeline/process node."""
        # Circle
        size = 0.35 if is_active else 0.25
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - size/2), Inches(y - size/2),
            Inches(size), Inches(size)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = primary_color if is_active else secondary_color
        circle.line.fill.background()
        
        # Label below
        label = slide.shapes.add_textbox(
            Inches(x - 0.8), Inches(y + 0.3),
            Inches(1.6), Inches(0.5)
        )
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = text[:20]
        p.font.size = Pt(10)
        p.font.bold = is_active
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        return circle
    
    @staticmethod
    def create_horizontal_line(slide, x: float, y: float, width: float, 
                              color, thickness: float = 2):
        """Create a horizontal line."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Pt(thickness)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    @staticmethod
    def create_vertical_bar(slide, x: float, y: float, height: float, 
                           color, width: float = 0.08):
        """Create a vertical accent bar."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape


class SlideRendererPro:
    """
    Professional slide renderer for meeting-ready presentations.
    
    Features:
    - Multiple layout styles per intent
    - Card-based layouts
    - Icon integration
    - Statistics/metrics displays
    - Timeline/process views
    - Professional typography
    """
    
    WIDTH = 13.333
    HEIGHT = 7.5
    MARGIN = 0.5
    
    def __init__(self, theme: ThemeColorScheme):
        self.theme = theme
        self.overflow = TextOverflowEngine()
        self.factory = ProShapeFactory()
        self.slide_count = 0
    
    def get_layout(self, slide_type: str) -> Dict[str, BoundingBox]:
        """Get layout for metrics evaluation."""
        w = self.WIDTH - 2 * self.MARGIN
        return {
            'title': BoundingBox(self.MARGIN, 0.4, w, 1.0),
            'content': BoundingBox(self.MARGIN, 1.5, w, 5.5)
        }
    
    def render(self, prs: Presentation, data: Dict) -> Any:
        """Render a slide with professional styling."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_type = data.get('slide_type', 'content')
        intent = data.get('intent', '')
        
        self.slide_count += 1
        
        # Route to appropriate renderer based on intent first, then slide_type
        # Only render as hero slide for actual title slides, not vision/cover intents after first slide
        if slide_type == 'title' and self.slide_count == 1:
            return self._render_hero_slide(slide, data)
        elif intent == 'closing' or slide_type == 'closing':
            return self._render_closing_slide(slide, data)
        elif intent == 'comparison' or slide_type == 'two_column' or slide_type == 'comparison':
            return self._render_comparison_slide(slide, data)
        elif intent == 'agenda':
            return self._render_agenda_slide(slide, data)
        elif intent == 'data_insight':
            return self._render_metrics_slide(slide, data)
        elif intent == 'framework':
            return self._render_process_slide(slide, data)
        elif intent == 'case_study':
            return self._render_case_slide(slide, data)
        elif intent in ['summary', 'call_to_action']:
            return self._render_closing_slide(slide, data)
        elif slide_type == 'section':
            return self._render_section_slide(slide, data)
        else:
            # Alternate between card and standard layouts for content slides
            if self.slide_count % 2 == 0:
                return self._render_card_layout(slide, data)
            else:
                return self._render_standard_layout(slide, data)
    
    def _render_hero_slide(self, slide, data: Dict):
        """Render an impactful hero/cover slide - matching closing slide style."""
        w = self.WIDTH
        h = self.HEIGHT
        
        # Full background - same style as closing slide
        self.factory.create_card(
            slide, BoundingBox(0, 0, w, h),
            self.theme.get_rgb('primary')
        )
        
        # Accent border at top
        self.factory.create_horizontal_line(
            slide, 0, 0, w,
            self.theme.get_rgb('accent'),
            thickness=10
        )
        
        # Decorative circles - same as closing slide
        self.factory.create_icon_circle(
            slide, -0.8, h - 2.2, '',
            self._lighten(self.theme.get_rgb('secondary'), 0.2),
            self.theme.get_rgb('secondary'),
            size=2.5
        )
        self.factory.create_icon_circle(
            slide, w - 1.5, -0.5, '',
            self._lighten(self.theme.get_rgb('accent'), 0.3),
            self.theme.get_rgb('accent'),
            size=2
        )
        
        # Title with smart typography - centered like closing
        title = data.get('title', '')
        if title:
            title_box = BoundingBox(self.MARGIN, 2.2, w - 2*self.MARGIN, 1.8)
            title_result = smart_typography.fit_text_smart(
                title, title_box,
                base_font_size=44,
                min_font_size=28,
                max_font_size=52
            )
            self.factory.create_text_box(
                slide,
                title_box,
                title_result['text'],
                title_result['font_size'],
                self.theme.get_rgb('text_light'),
                bold=True,
                align=PP_ALIGN.CENTER,
                font_name=Typography.HEADING
            )
        
        # Subtitle - elegant italic below title
        subtitle = data.get('subtitle', '')
        if subtitle:
            subtitle_box = BoundingBox(self.MARGIN + 1, 4.2, w - 2*self.MARGIN - 2, 1.0)
            subtitle_result = smart_typography.fit_text_smart(
                subtitle, subtitle_box,
                base_font_size=22,
                min_font_size=14,
                max_font_size=26
            )
            self.factory.create_text_box(
                slide,
                subtitle_box,
                subtitle_result['text'],
                subtitle_result['font_size'],
                self._lighten(self.theme.get_rgb('text_light'), 0.2),
                italic=True,
                align=PP_ALIGN.CENTER,
                font_name=Typography.BODY
            )
        
        # Bottom accent bar
        self.factory.create_card(
            slide, BoundingBox(w/2 - 1, h - 0.8, 2, 0.1),
            self.theme.get_rgb('accent')
        )
        
        return slide
    
    def _render_section_slide(self, slide, data: Dict):
        """Render a section divider with visual impact."""
        w = self.WIDTH
        h = self.HEIGHT
        
        # Large accent circle (background)
        self.factory.create_icon_circle(
            slide, w - 3, h/2 - 1.5, '',
            self._lighten(self.theme.get_rgb('accent'), 0.85),
            self.theme.get_rgb('accent'),
            size=3
        )
        
        # Vertical accent bar
        self.factory.create_vertical_bar(
            slide, 0, 0, h,
            self.theme.get_rgb('primary'),
            width=0.12
        )
        
        # Section number
        self.factory.create_text_box(
            slide,
            BoundingBox(0.5, 1.5, 1.5, 1),
            f"0{self.slide_count}",
            48,
            self._lighten(self.theme.get_rgb('primary'), 0.6),
            bold=True
        )
        
        # Title with smart typography
        title = data.get('title', '')
        if title:
            title_box = BoundingBox(0.5, 2.8, w - 4, 1.6)
            title_result = smart_typography.fit_text_smart(
                title, title_box,
                base_font_size=38,
                min_font_size=26,
                max_font_size=44
            )
            self.factory.create_text_box(
                slide,
                title_box,
                title_result['text'],
                title_result['font_size'],
                self.theme.get_rgb('primary'),
                bold=True,
                align=PP_ALIGN.LEFT
            )
        
        # Subtitle with smart typography
        subtitle = data.get('subtitle', '')
        if subtitle:
            subtitle_box = BoundingBox(0.5, 4.5, w - 4, 1.0)
            subtitle_result = smart_typography.fit_text_smart(
                subtitle, subtitle_box,
                base_font_size=18,
                min_font_size=14,
                max_font_size=22
            )
            self.factory.create_text_box(
                slide,
                subtitle_box,
                subtitle_result['text'],
                subtitle_result['font_size'],
                self.theme.get_rgb('text_secondary') if hasattr(self.theme, 'text_secondary') else self._lighten(self.theme.get_rgb('text_dark'), 0.4),
                italic=True
            )
        
        return slide
    
    def _render_standard_layout(self, slide, data: Dict):
        """Render standard content slide with professional styling and smart typography."""
        w = self.WIDTH - 2 * self.MARGIN
        
        # Top accent line
        self.factory.create_horizontal_line(
            slide, 0, 0, self.WIDTH,
            self.theme.get_rgb('primary'),
            thickness=4
        )
        
        # Title with smart typography
        title = data.get('title', '')
        if title:
            title_box = BoundingBox(self.MARGIN, 0.25, w, 0.9)
            title_result = smart_typography.fit_text_smart(
                title, title_box,
                base_font_size=26,
                min_font_size=18,
                max_font_size=30
            )
            self.factory.create_text_box(
                slide,
                title_box,
                title_result['text'],
                title_result['font_size'],
                self.theme.get_rgb('primary'),
                bold=True,
                font_name=Typography.HEADING
            )
        
        # Accent underline
        self.factory.create_horizontal_line(
            slide, self.MARGIN, 1.2, 2.5,
            self.theme.get_rgb('accent'),
            thickness=4
        )
        
        # Body points with smart sizing based on content amount
        points = data.get('body_points') or []
        if points:
            processed = self._process_points(points, max_items=6)
            # Calculate adaptive font size based on number of points
            num_points = len(processed)
            if num_points <= 3:
                font_size = 20  # Larger font for fewer items
            elif num_points <= 4:
                font_size = 18
            elif num_points <= 5:
                font_size = 16
            else:
                font_size = 15  # Smaller for many items
            
            self.factory.create_styled_bullet_list(
                slide,
                BoundingBox(self.MARGIN, 1.5, w, 5.5),
                processed,
                font_size,
                self.theme.get_rgb('text_dark'),
                self.theme.get_rgb('accent')
            )
        
        # Slide number
        self.factory.create_text_box(
            slide,
            BoundingBox(self.WIDTH - 1, self.HEIGHT - 0.5, 0.8, 0.4),
            str(self.slide_count),
            10,
            self._lighten(self.theme.get_rgb('text_dark'), 0.6),
            align=PP_ALIGN.RIGHT
        )
        
        return slide
    
    def _render_card_layout(self, slide, data: Dict):
        """Render content in card-based layout with vector icons."""
        w = self.WIDTH - 2 * self.MARGIN
        intent = data.get('intent', '')
        
        # Header bar
        self.factory.create_card(
            slide, BoundingBox(0, 0, self.WIDTH, 1.3),
            self.theme.get_rgb('primary')
        )
        
        # Title with smart typography
        title = data.get('title', '')
        if title:
            title_box = BoundingBox(self.MARGIN, 0.25, w, 0.85)
            title_result = smart_typography.fit_text_smart(
                title, title_box,
                base_font_size=26,
                min_font_size=18,
                max_font_size=30
            )
            self.factory.create_text_box(
                slide,
                title_box,
                title_result['text'],
                title_result['font_size'],
                self.theme.get_rgb('text_light'),
                bold=True,
                font_name=Typography.HEADING
            )
        
        # Content cards with vector icons
        points = data.get('body_points') or []
        if points:
            processed = self._process_points(points, max_items=4)
            card_width = (w - 0.4) / 2
            card_height = 2.4
            
            # Icon creators for each card
            icon_types = ['lightbulb', 'chart_bars', 'target', 'arrow_up']
            icon_colors = [
                self.theme.get_rgb('accent'),
                self.theme.get_rgb('secondary'),
                self.theme.get_rgb('primary'),
                self._lighten(self.theme.get_rgb('accent'), 0.2)
            ]
            
            for i, point in enumerate(processed[:4]):
                row = i // 2
                col = i % 2
                x = self.MARGIN + col * (card_width + 0.4)
                y = 1.7 + row * (card_height + 0.3)
                
                # Card background
                self.factory.create_card(
                    slide,
                    BoundingBox(x, y, card_width, card_height),
                    self._lighten(self.theme.get_rgb('secondary'), 0.92),
                    border_color=self._lighten(self.theme.get_rgb('secondary'), 0.6)
                )
                
                # Card text with smart typography
                text = point.get('text', '')
                card_text_box = BoundingBox(x + 0.3, y + 0.4, card_width - 0.6, card_height - 0.8)
                card_text_result = smart_typography.fit_text_smart(
                    text, card_text_box,
                    base_font_size=16,
                    min_font_size=12,
                    max_font_size=18
                )
                self.factory.create_text_box(
                    slide,
                    card_text_box,
                    card_text_result['text'],
                    card_text_result['font_size'],
                    self.theme.get_rgb('text_dark'),
                    align=PP_ALIGN.LEFT,
                    line_spacing=1.4,
                    font_name=Typography.BODY
                )
        
        return slide
    
    def _render_comparison_slide(self, slide, data: Dict):
        """Render professional comparison/two-column slide."""
        w = self.WIDTH - 2 * self.MARGIN
        col_width = (w - 0.6) / 2
        
        # Header
        self.factory.create_card(
            slide, BoundingBox(0, 0, self.WIDTH, 1.4),
            self.theme.get_rgb('primary')
        )
        
        # Title with smart typography - left aligned to avoid VS overlap
        title = data.get('title', '')
        if title:
            title_box = BoundingBox(self.MARGIN, 0.4, w * 0.7, 0.7)
            title_result = smart_typography.fit_text_smart(
                title, title_box,
                base_font_size=24,
                min_font_size=18,
                max_font_size=28
            )
            self.factory.create_text_box(
                slide,
                title_box,
                title_result['text'],
                title_result['font_size'],
                self.theme.get_rgb('text_light'),
                bold=True,
                align=PP_ALIGN.LEFT
            )
        
        # Left column card
        left_header = data.get('left_header') or 'Option A'
        self.factory.create_card(
            slide,
            BoundingBox(self.MARGIN, 1.7, col_width, 5.3),
            self._lighten(self.theme.get_rgb('secondary'), 0.85)
        )
        self.factory.create_card(
            slide,
            BoundingBox(self.MARGIN, 1.7, col_width, 0.6),
            self._lighten(self.theme.get_rgb('secondary'), 0.4)
        )
        self.factory.create_text_box(
            slide,
            BoundingBox(self.MARGIN + 0.2, 1.8, col_width - 0.4, 0.4),
            f"{left_header[:30]}",
            14,
            self.theme.get_rgb('primary'),
            bold=True,
            align=PP_ALIGN.CENTER
        )
        
        # Left content with adaptive font sizing
        left_items = data.get('left_column') or []
        if left_items:
            processed = self._process_column_items(left_items)
            num_items = len(processed)
            # Calculate font size based on number of items
            font_size = 16 if num_items <= 3 else (14 if num_items <= 5 else 12)
            item_height = min(0.95, 4.0 / max(num_items, 1))
            y_pos = 2.5
            for item in processed:
                item_text = f"• {item['text']}"
                item_box = BoundingBox(self.MARGIN + 0.3, y_pos, col_width - 0.6, item_height)
                item_result = smart_typography.fit_text_smart(
                    item_text, item_box,
                    base_font_size=font_size,
                    min_font_size=11,
                    max_font_size=18
                )
                self.factory.create_text_box(
                    slide,
                    item_box,
                    item_result['text'],
                    item_result['font_size'],
                    self.theme.get_rgb('text_dark')
                )
                y_pos += item_height + 0.1
        
        # Right column card
        right_header = data.get('right_header') or 'Option B'
        self.factory.create_card(
            slide,
            BoundingBox(self.MARGIN + col_width + 0.6, 1.7, col_width, 5.3),
            self._lighten(self.theme.get_rgb('accent'), 0.85)
        )
        self.factory.create_card(
            slide,
            BoundingBox(self.MARGIN + col_width + 0.6, 1.7, col_width, 0.6),
            self._lighten(self.theme.get_rgb('accent'), 0.4)
        )
        self.factory.create_text_box(
            slide,
            BoundingBox(self.MARGIN + col_width + 0.8, 1.8, col_width - 0.4, 0.4),
            f"{right_header[:30]}",
            14,
            self.theme.get_rgb('primary'),
            bold=True,
            align=PP_ALIGN.CENTER
        )
        
        # Right content with adaptive font sizing
        right_items = data.get('right_column') or []
        if right_items:
            processed = self._process_column_items(right_items)
            num_items = len(processed)
            font_size = 16 if num_items <= 3 else (14 if num_items <= 5 else 12)
            item_height = min(0.95, 4.0 / max(num_items, 1))
            y_pos = 2.5
            for item in processed:
                item_text = f"• {item['text']}"
                item_box = BoundingBox(self.MARGIN + col_width + 0.9, y_pos, col_width - 0.6, item_height)
                item_result = smart_typography.fit_text_smart(
                    item_text, item_box,
                    base_font_size=font_size,
                    min_font_size=11,
                    max_font_size=18
                )
                self.factory.create_text_box(
                    slide,
                    item_box,
                    item_result['text'],
                    item_result['font_size'],
                    self.theme.get_rgb('text_dark')
                )
                y_pos += item_height + 0.1
        
        return slide
    
    def _render_metrics_slide(self, slide, data: Dict):
        """Render data/metrics focused slide with stat boxes."""
        w = self.WIDTH - 2 * self.MARGIN
        
        # Header with line
        self.factory.create_horizontal_line(
            slide, 0, 0, self.WIDTH,
            self.theme.get_rgb('primary'),
            thickness=4
        )
        
        # Title with smart typography
        title = data.get('title', '')
        if title:
            title_box = BoundingBox(self.MARGIN, 0.25, w, 0.8)
            title_result = smart_typography.fit_text_smart(
                title, title_box,
                base_font_size=26,
                min_font_size=18,
                max_font_size=30
            )
            self.factory.create_text_box(
                slide,
                title_box,
                title_result['text'],
                title_result['font_size'],
                self.theme.get_rgb('primary'),
                bold=True,
                font_name=Typography.HEADING
            )
        
        # Metrics boxes - larger and vertically centered
        points = data.get('body_points') or []
        if points:
            num_boxes = min(len(points), 4)
            # Adaptive sizing based on number of boxes
            if num_boxes <= 2:
                box_width = 5.0
                box_height = 3.2
                gap = 1.0
            elif num_boxes == 3:
                box_width = 3.8
                box_height = 3.0
                gap = 0.6
            else:
                box_width = 3.0
                box_height = 2.8
                gap = 0.5
            
            total_width = num_boxes * box_width + (num_boxes - 1) * gap
            start_x = (self.WIDTH - total_width) / 2
            # Vertically center the boxes (slide height ~7.5, title takes ~1.2)
            available_height = self.HEIGHT - 1.5  # 6.0 inches available
            start_y = 1.5 + (available_height - box_height) / 2  # Center vertically
            
            colors = [
                self.theme.get_rgb('primary'),
                self.theme.get_rgb('secondary'),
                self.theme.get_rgb('accent'),
                self._lighten(self.theme.get_rgb('primary'), 0.3)
            ]
            
            for i, point in enumerate(points[:4]):
                text = point.get('text', '') if isinstance(point, dict) else str(point)
                # Try to extract number from text (format: "number: label" or just text)
                parts = text.split(':')
                if len(parts) >= 2:
                    number = parts[0].strip()
                    label = ':'.join(parts[1:]).strip()
                else:
                    number = f"#{i+1}"
                    label = text
                
                x = start_x + i * (box_width + gap)
                self.factory.create_stat_box(
                    slide, x, start_y, box_width, box_height,
                    number, label,
                    self._lighten(colors[i], 0.85),
                    colors[i],
                    self.theme.get_rgb('text_dark')
                )
            
            # Supporting text below metrics - if there's room
            if len(points) > 4:
                remaining = points[4:]
                y = start_y + box_height + 0.8
                for p in remaining[:2]:
                    text = p.get('text', '') if isinstance(p, dict) else str(p)
                    self.factory.create_text_box(
                        slide,
                        BoundingBox(self.MARGIN, y, w, 0.6),
                        text[:100],
                        14,
                        self.theme.get_rgb('text_dark'),
                        align=PP_ALIGN.CENTER
                    )
                    y += 0.8
        
        return slide
    
    def _render_process_slide(self, slide, data: Dict):
        """Render framework/process slide with timeline."""
        w = self.WIDTH - 2 * self.MARGIN
        
        # Header
        self.factory.create_card(
            slide, BoundingBox(0, 0, self.WIDTH, 1.2),
            self.theme.get_rgb('primary')
        )
        
        title = data.get('title', '')
        if title:
            title_box = BoundingBox(self.MARGIN, 0.25, w, 0.7)
            title_result = smart_typography.fit_text_smart(
                f"⚙️ {title}", title_box,
                base_font_size=22,
                min_font_size=16,
                max_font_size=26
            )
            self.factory.create_text_box(
                slide,
                title_box,
                title_result['text'],
                title_result['font_size'],
                self.theme.get_rgb('text_light'),
                bold=True
            )
        
        # Process steps
        points = data.get('body_points') or []
        if points:
            num_steps = min(len(points), 5)
            step_width = (w - 1) / num_steps
            y_line = 2.2
            
            # Connecting line
            self.factory.create_horizontal_line(
                slide, self.MARGIN + 0.5, y_line,
                w - 1,
                self._lighten(self.theme.get_rgb('secondary'), 0.5),
                thickness=3
            )
            
            # Calculate adaptive font size for steps
            step_font_size = 14 if num_steps <= 3 else (12 if num_steps <= 4 else 11)
            
            # Step nodes
            for i, point in enumerate(points[:5]):
                text = point.get('text', '') if isinstance(point, dict) else str(point)
                x = self.MARGIN + 0.5 + i * step_width + step_width/2
                
                # Node circle
                self.factory.create_icon_circle(
                    slide, x - 0.25, y_line - 0.25, str(i + 1),
                    self.theme.get_rgb('accent'),
                    self.theme.get_rgb('text_light'),
                    size=0.5
                )
                
                # Step card
                self.factory.create_card(
                    slide,
                    BoundingBox(x - step_width/2 + 0.1, 2.8, step_width - 0.2, 3.5),
                    self._lighten(self.theme.get_rgb('secondary'), 0.9)
                )
                
                # Step text with smart typography
                step_box = BoundingBox(x - step_width/2 + 0.15, 3.0, step_width - 0.3, 3.0)
                step_result = smart_typography.fit_text_smart(
                    text, step_box,
                    base_font_size=step_font_size,
                    min_font_size=9,
                    max_font_size=16
                )
                self.factory.create_text_box(
                    slide,
                    step_box,
                    step_result['text'],
                    step_result['font_size'],
                    self.theme.get_rgb('text_dark'),
                    align=PP_ALIGN.CENTER
                )
        
        return slide
    
    def _render_case_slide(self, slide, data: Dict):
        """Render case study/example slide."""
        w = self.WIDTH - 2 * self.MARGIN
        
        # Side accent
        self.factory.create_vertical_bar(
            slide, 0, 0, self.HEIGHT,
            self.theme.get_rgb('accent'),
            width=0.1
        )
        
        # Title with smart typography
        title = data.get('title', '')
        if title:
            title_box = BoundingBox(self.MARGIN, 0.4, w, 1.0)
            title_result = smart_typography.fit_text_smart(
                title, title_box,
                base_font_size=28,
                min_font_size=20,
                max_font_size=32
            )
            self.factory.create_text_box(
                slide,
                title_box,
                title_result['text'],
                title_result['font_size'],
                self.theme.get_rgb('primary'),
                bold=True
            )
        
        # Content with quote style and smart typography
        points = data.get('body_points') or []
        if points:
            # Quote bar
            self.factory.create_vertical_bar(
                slide, self.MARGIN, 1.7, 5.0,
                self.theme.get_rgb('accent'),
                width=0.06
            )
            
            processed = self._process_points(points, max_items=5)
            num_points = len(processed)
            point_height = min(1.1, 5.0 / max(num_points, 1))
            font_size = 18 if num_points <= 3 else (16 if num_points <= 4 else 14)
            y = 1.8
            for point in processed:
                point_box = BoundingBox(self.MARGIN + 0.3, y, w - 0.3, point_height)
                point_result = smart_typography.fit_text_smart(
                    point['text'], point_box,
                    base_font_size=font_size,
                    min_font_size=12,
                    max_font_size=20
                )
                self.factory.create_text_box(
                    slide,
                    point_box,
                    point_result['text'],
                    point_result['font_size'],
                    self.theme.get_rgb('text_dark')
                )
                y += point_height + 0.05
        
        return slide
    
    def _render_agenda_slide(self, slide, data: Dict):
        """Render agenda with numbered items."""
        w = self.WIDTH - 2 * self.MARGIN
        
        # Header
        self.factory.create_card(
            slide, BoundingBox(0, 0, self.WIDTH, 1.3),
            self.theme.get_rgb('primary')
        )
        
        title = data.get('title', 'Agenda')
        title_box = BoundingBox(self.MARGIN, 0.3, w, 0.7)
        title_result = smart_typography.fit_text_smart(
            title, title_box,
            base_font_size=24,
            min_font_size=18,
            max_font_size=28
        )
        self.factory.create_text_box(
            slide,
            title_box,
            title_result['text'],
            title_result['font_size'],
            self.theme.get_rgb('text_light'),
            bold=True
        )
        
        # Numbered items with smart typography
        points = data.get('body_points') or []
        num_items = min(len(points), 6)
        y_start = 1.7
        # Calculate item height based on number of items
        item_height = min(1.0, 4.8 / max(num_items, 1))
        font_size = 18 if num_items <= 4 else (16 if num_items <= 5 else 14)
        
        for i, point in enumerate(points[:6]):
            text = point.get('text', '') if isinstance(point, dict) else str(point)
            
            # Number circle
            self.factory.create_icon_circle(
                slide, self.MARGIN, y_start + i * item_height, str(i + 1),
                self.theme.get_rgb('accent'),
                self.theme.get_rgb('text_light'),
                size=0.45
            )
            
            # Item text with smart typography
            item_box = BoundingBox(self.MARGIN + 0.65, y_start + i * item_height, w - 1, item_height - 0.15)
            item_result = smart_typography.fit_text_smart(
                text, item_box,
                base_font_size=font_size,
                min_font_size=12,
                max_font_size=20
            )
            self.factory.create_text_box(
                slide,
                item_box,
                item_result['text'],
                item_result['font_size'],
                self.theme.get_rgb('text_dark')
            )
            
            # Separator line
            if i < num_items - 1:
                self.factory.create_horizontal_line(
                    slide, self.MARGIN + 0.65,
                    y_start + (i + 1) * item_height - 0.1,
                    w - 1,
                    self._lighten(self.theme.get_rgb('text_dark'), 0.85),
                    thickness=1
                )
        
        return slide
    
    def _render_closing_slide(self, slide, data: Dict):
        """Render professional closing slide with smart typography."""
        w = self.WIDTH
        h = self.HEIGHT
        
        # Full background
        self.factory.create_card(
            slide, BoundingBox(0, 0, w, h),
            self.theme.get_rgb('primary')
        )
        
        # Decorative elements
        self.factory.create_icon_circle(
            slide, -0.8, h - 2.2, '',
            self._lighten(self.theme.get_rgb('secondary'), 0.2),
            self.theme.get_rgb('secondary'),
            size=2.5
        )
        self.factory.create_icon_circle(
            slide, w - 1.5, -0.5, '',
            self._lighten(self.theme.get_rgb('accent'), 0.3),
            self.theme.get_rgb('accent'),
            size=2
        )
        
        # Title - use smart typography to fit
        title = data.get('title', 'Thank You')
        title_box = BoundingBox(self.MARGIN, 2.0, w - 2*self.MARGIN, 1.5)
        title_result = smart_typography.fit_text_smart(
            title, title_box,
            base_font_size=44,
            min_font_size=28,
            max_font_size=52
        )
        self.factory.create_text_box(
            slide,
            title_box,
            title_result['text'],
            title_result['font_size'],
            self.theme.get_rgb('text_light'),
            bold=True,
            align=PP_ALIGN.CENTER,
            font_name=Typography.HEADING
        )
        
        # Subtitle - use smart typography
        subtitle = data.get('subtitle', '')
        if subtitle:
            subtitle_box = BoundingBox(self.MARGIN + 1, 4.3, w - 2*self.MARGIN - 2, 1.0)
            subtitle_result = smart_typography.fit_text_smart(
                subtitle, subtitle_box,
                base_font_size=22,
                min_font_size=14,
                max_font_size=26
            )
            self.factory.create_text_box(
                slide,
                subtitle_box,
                subtitle_result['text'],
                subtitle_result['font_size'],
                self._lighten(self.theme.get_rgb('text_light'), 0.15),
                align=PP_ALIGN.CENTER
            )
        
        # Bottom accent
        self.factory.create_horizontal_line(
            slide, w/2 - 1.5, 5.2, 3,
            self.theme.get_rgb('accent'),
            thickness=4
        )
        
        return slide
    
    def _process_points(self, points: List, max_items: int = 5) -> List[Dict]:
        """Process and limit body points."""
        processed = []
        for p in points[:max_items]:
            if isinstance(p, dict):
                processed.append({
                    'text': p.get('text', '')[:100],
                    'level': p.get('level', 0),
                    'priority': p.get('priority', 'normal')
                })
            else:
                processed.append({'text': str(p)[:100], 'level': 0, 'priority': 'normal'})
        return processed
    
    def _process_column_items(self, items: List) -> List[Dict]:
        """Process column items."""
        processed = []
        for item in items[:4]:
            if isinstance(item, dict):
                processed.append({'text': item.get('text', '')[:60], 'level': 0})
            else:
                processed.append({'text': str(item)[:60], 'level': 0})
        return processed
    
    def _create_vector_icon(self, slide, icon_type: str, x: float, y: float, 
                            size: float, color: RGBColor):
        """Create a vector icon based on type."""
        colors = [color, self._lighten(color, 0.3), self._lighten(color, 0.5)]
        
        if icon_type == 'lightbulb':
            VectorIcons.create_lightbulb(slide, x, y, size, color)
        elif icon_type == 'chart_bars':
            VectorIcons.create_chart_bars(slide, x, y, size, colors)
        elif icon_type == 'target':
            VectorIcons.create_target(slide, x, y, size, colors)
        elif icon_type == 'gear':
            VectorIcons.create_gear(slide, x, y, size, color)
        elif icon_type == 'arrow_up':
            VectorIcons.create_arrow_up(slide, x, y, size, color)
        elif icon_type == 'checkmark':
            VectorIcons.create_checkmark(slide, x, y, size, color)
        elif icon_type == 'document':
            VectorIcons.create_document(slide, x, y, size, color)
        elif icon_type == 'users':
            VectorIcons.create_users(slide, x, y, size, color)
        else:
            # Default: simple colored circle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y),
                Inches(size), Inches(size)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
    
    def _lighten(self, rgb_color: RGBColor, factor: float) -> RGBColor:
        """Lighten a color."""
        r = int(rgb_color[0] + (255 - rgb_color[0]) * factor)
        g = int(rgb_color[1] + (255 - rgb_color[1]) * factor)
        b = int(rgb_color[2] + (255 - rgb_color[2]) * factor)
        return RGBColor(min(r, 255), min(g, 255), min(b, 255))
    
    def _darken(self, rgb_color: RGBColor, factor: float) -> RGBColor:
        """Darken a color."""
        r = int(rgb_color[0] * (1 - factor))
        g = int(rgb_color[1] * (1 - factor))
        b = int(rgb_color[2] * (1 - factor))
        return RGBColor(max(r, 0), max(g, 0), max(b, 0))


__all__ = ['ProShapeFactory', 'SlideRendererPro', 'VectorIcons', 'Typography']

